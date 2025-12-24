# modules/doc_processor.py
import os
from docx import Document
from langchain_community.document_loaders import PyPDFLoader, TextLoader, UnstructuredWordDocumentLoader
from .translator import BailianTranslator
import config

# 新增：只提取文本，不总结
def extract_text_from_file(file_path: str, max_chars: int = 1000) -> str:
    """
    读取文件内容并返回字符串
    max_chars: 限制读取长度，防止Token爆炸
    """
    print(f"[DEBUG] 正在尝试读取文件: {file_path}")
    print(f"[DEBUG] 文件路径是否存在: {os.path.exists(file_path)}")

    if not os.path.exists(file_path):
        print(f"[DEBUG] 文件不存在，返回: [文件不存在]")
        return "[文件不存在]"

    ext = os.path.splitext(file_path)[1].lower()
    print(f"[DEBUG] 文件扩展名: {ext}")

    try:
        if ext == ".pdf":
            print(f"[DEBUG] 使用 PyPDFLoader 读取 PDF 文件")
            loader = PyPDFLoader(file_path)
            print(f"[DEBUG] 开始加载文档...")
            docs = loader.load()
            print(f"[DEBUG] 加载了 {len(docs)} 个文档片段")
            full_text = "\n".join([d.page_content for d in docs])
            print(f"[DEBUG] 合并后的文本长度: {len(full_text)}")
        elif ext == ".txt":
            print(f"[DEBUG] 使用 TextLoader 读取 TXT 文件")
            loader = TextLoader(file_path, encoding="utf-8")
            print(f"[DEBUG] 开始加载文档...")
            docs = loader.load()
            print(f"[DEBUG] 加载了 {len(docs)} 个文档片段")
            full_text = "\n".join([d.page_content for d in docs])
            print(f"[DEBUG] 合并后的文本长度: {len(full_text)}")
        elif ext in [".docx", ".doc"]:
            print(f"[DEBUG] 使用 UnstructuredWordDocumentLoader 读取 Word 文件")
            loader = UnstructuredWordDocumentLoader(file_path)
            print(f"[DEBUG] 开始加载文档...")
            docs = loader.load()
            print(f"[DEBUG] 加载了 {len(docs)} 个文档片段")
            full_text = "\n".join([d.page_content for d in docs])
            print(f"[DEBUG] 合并后的文本长度: {len(full_text)}")
        elif ext == ".pptx":
            print(f"[DEBUG] 使用 PowerPoint 文件处理器读取 PPTX 文件")
            # 导入 PowerPoint 处理模块
            from pptx import Presentation
            # 直接读取 PPTX 文件内容
            prs = Presentation(file_path)
            full_text = []
            for slide_num, slide in enumerate(prs.slides):
                slide_text = f"\n--- 幻面 {slide_num + 1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text += shape.text + "\n"
                if slide_text.strip():  # 如果幻灯片有文本内容
                    full_text.append(slide_text)
            full_text = "".join(full_text)
            print(f"[DEBUG] 从 PPTX 提取的文本长度: {len(full_text)}")
        else:
            result = f"[不支持的文件格式: {ext}]"
            print(f"[DEBUG] 不支持的文件格式，返回: {result}")
            return result

        # 简单清洗空格
        full_text = full_text.strip()
        print(f"[DEBUG] 清洗后的文本长度: {len(full_text)}")

        if len(full_text) > max_chars:
            result = full_text[:max_chars] + "...(剩余内容已截断)"
            print(f"[DEBUG] 文本超过最大长度限制，截断后长度: {len(result)}")
            return result
        result = full_text if full_text else "[空文件]"
        print(f"[DEBUG] 最终返回文本长度: {len(result) if result != '[空文件]' else 0}, 内容预览: {result[:100]}...")
        return result

    except Exception as e:
        result = f"[读取文件出错: {str(e)}]"
        print(f"[DEBUG] 读取文件时发生异常: {e}")
        print(f"[DEBUG] 异常类型: {type(e).__name__}")
        import traceback
        print(f"[DEBUG] 异常堆栈: {traceback.format_exc()}")
        return result

# 原有函数保持不变，但可以调用上面的函数来简化逻辑
def process_document_summary(file_path: str, task_type: str = "summarize", target_lang: str = "Chinese"):
    print(f"[DEBUG_SUMMARY] 开始处理文档总结，文件路径: {file_path}, 任务类型: {task_type}")
    # 复用上面的提取逻辑
    full_text = extract_text_from_file(file_path, max_chars=15000)
    print(f"[DEBUG_SUMMARY] 提取的文本: {full_text[:200]}...")

    if full_text.startswith("["): # 简单的错误检查
        if "出错" in full_text or "不支持" in full_text:
            print(f"[DEBUG_SUMMARY] 检测到错误信息，直接返回: {full_text}")
            return full_text

    print(f"[DEBUG_SUMMARY] 调用翻译器进行{task_type}，目标语言: {target_lang}")
    translator = BailianTranslator(config.DASHSCOPE_API_KEY)
    result = translator._call_api(full_text, mode=task_type, target_lang=target_lang)
    print(f"[DEBUG_SUMMARY] 翻译结果: {result[:200]}...")
    return result

def save_text_to_docx(text: str, output_path: str):
    """
    将文本保存为 Word 文档
    """
    doc = Document()
    # 按换行符分割段落，避免所有文字挤在一段
    paragraphs = text.split('\n')
    for p in paragraphs:
        if p.strip():
            doc.add_paragraph(p)
            
    doc.save(output_path)
    return output_path
    